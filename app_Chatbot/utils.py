import os
from .models import Arquivo
import logging
import PyPDF2
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from dotenv import load_dotenv
import os
import pickle
from sklearn.metrics.pairwise import cosine_similarity
from .models import Conteudo, Arquivo

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Carrega as variáveis de ambiente do arquivo .env
load_dotenv()
os.getenv("")

# Função para verificar se o arquivo já existe
def file_exists(file_name):
    # Verifica se o arquivo existe na pasta de uploads
    file_path = os.path.join('uploads/', file_name)
    if os.path.exists(file_path):
        return True
    
    # Verifica se o arquivo existe no banco de dados
    if Arquivo.objects.filter(nome_do_arquivo=file_name).exists():
        return True
    
    return False

# Função para extrair texto de arquivos PDF
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PyPDF2.PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Função para extrair texto de arquivos DOCX (Word)
def get_docx_text(docx_files):
    text = ""
    for docx_file in docx_files:
        doc = DocxDocument(docx_file)
        for paragraph in doc.paragraphs:
            text += paragraph.text
    return text

# Função para extrair texto de arquivos XLSX (Excel)
def get_xlsx_text(xlsx_files):
    text = ""
    for xlsx_file in xlsx_files:
        wb = load_workbook(filename=xlsx_file)
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + " "
    return text

# Função para extrair texto de arquivos PPTX (PowerPoint)
def get_pptx_text(pptx_files):
    text = ""
    for pptx_file in pptx_files:
        prs = Presentation(pptx_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Função para dividir o texto em pedaços menores
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks

# Função para criar e salvar os vetores de armazenamento a partir dos pedaços de texto
def create_and_store_embeddings(text_chunks, arquivo):
    embeddings_model = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    # Salvando os chunks e embeddings no banco de dados
    for chunk in text_chunks:
        try:
            embedding = embeddings_model.embed_documents([chunk])[0]
            # Convertendo embedding para uma lista serializável
            embedding_list = list(embedding)
            serialized_embedding = pickle.dumps(embedding_list)  # Serializando o embedding
            Conteudo.objects.create(
                texto=chunk,
                id_arquivo_id=arquivo,  # Corrigido para o nome correto do campo
                embeddings=serialized_embedding
            )
            logger.debug(f"Stored chunk: {chunk}")
        except Exception as e:
            logger.error(f"Erro ao criar e armazenar embedding: {str(e)}")

# Função para criar a cadeia de conversação
def get_conversational_chain():
    prompt_template = """
    Responda à pergunta o mais detalhadamente possível, usando o contexto fornecido. 
    Certifique-se de fornecer todos os detalhes necessários. 
    Se a resposta não estiver no contexto fornecido, apenas diga: "resposta não disponível no contexto", 
    e evite fornecer uma resposta incorreta.\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Responder:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

class TextDocument:
    def __init__(self, page_content, metadata=None):
        self.page_content = page_content
        self.metadata = metadata or {}

# Função para processar a entrada do usuário
def user_input(user_question, conversation):
    embeddings_model = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    question_embedding = embeddings_model.embed_documents([user_question])[0]
    
    # Carregar todos os conteúdos e seus embeddings do banco de dados
    conteudos = Conteudo.objects.all()
    chunks = []
    for conteudo in conteudos:
        embedding = pickle.loads(conteudo.embeddings)
        chunks.append((conteudo.texto, embedding))
    
    # Realizar a busca por similaridade usando embeddings
    docs = []
    for chunk_text, embedding in chunks:
        similarity = cosine_similarity([question_embedding], [embedding])[0][0]
        if similarity > 0.5:  # Threshold de similaridade, ajustável
            docs.append(TextDocument(page_content=chunk_text))
    
    chain = get_conversational_chain()
    response = chain.invoke({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    conversation.append({"user_question": user_question, "bot_response": response["output_text"]})
    return response["output_text"], conversation
